#!/usr/bin/env python3
"""Build a premium CiteVision architectures PowerPoint from the pattern map."""
from __future__ import annotations

import base64
import io
import zlib
from pathlib import Path

import requests
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "CiteVision-Architectures.pptx"
ASSETS = ROOT / "_pptx_assets"
ASSETS.mkdir(exist_ok=True)

# Palette aligned with docs/architectures/index.html
BG = RGBColor(0x0B, 0x10, 0x12)
BG2 = RGBColor(0x12, 0x18, 0x1A)
INK = RGBColor(0xE8, 0xEF, 0xE9)
MUTED = RGBColor(0x9A, 0xAB, 0xA3)
ACCENT = RGBColor(0xD4, 0xA5, 0x74)
TEAL = RGBColor(0x4A, 0x9B, 0x84)
OK = RGBColor(0x7C, 0xBC, 0x8F)

W, H = Inches(13.333), Inches(7.5)  # 16:9


def set_run(run, *, size=18, bold=False, color=INK, font="Calibri"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def fill_solid(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_bg(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    fill_solid(shape, BG)
    # subtle top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.06))
    fill_solid(bar, TEAL)
    # copper accent corner
    corner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(12.6), Inches(7.15), Inches(0.55), Inches(0.22)
    )
    fill_solid(corner, ACCENT)


def add_textbox(slide, left, top, width, height, text, *, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font)
    return box


def add_para(tf, text, *, size=14, bold=False, color=MUTED, space_before=6):
    p = tf.add_paragraph()
    p.space_before = Pt(space_before)
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return p


def mermaid_to_png(diagram: str, name: str) -> Path | None:
    """Fetch PNG from kroki.io (mermaid)."""
    dest = ASSETS / f"{name}.png"
    if dest.exists() and dest.stat().st_size > 800:
        return dest
    # Kroki: POST mermaid -> png
    try:
        r = requests.post(
            "https://kroki.io/mermaid/png",
            data=diagram.encode("utf-8"),
            headers={"Content-Type": "text/plain"},
            timeout=45,
        )
        if r.status_code == 200 and r.content[:8] == b"\x89PNG\r\n\x1a\n":
            dest.write_bytes(r.content)
            return dest
        # fallback mermaid.ink
        raw = base64.urlsafe_b64encode(diagram.encode("utf-8")).decode("ascii").rstrip("=")
        url = f"https://mermaid.ink/img/{raw}?type=png&bgColor=!0e1416"
        r2 = requests.get(url, timeout=45)
        if r2.status_code == 200 and len(r2.content) > 500:
            dest.write_bytes(r2.content)
            return dest
    except Exception as exc:
        print(f"[WARN] diagram {name}: {exc}")
    return None


SLIDES = [
    {
        "id": "intro",
        "kind": "title",
        "title": "CitéVision",
        "subtitle": "Cartographie des architectures métier & produit",
        "body": "Zone → IA → règle → preuve\n10 patterns de règles · 6 blocs produit\nPour reconnaître le mécanisme de n’importe quelle règle du catalogue.",
    },
    {
        "id": "read",
        "kind": "bullets",
        "eyebrow": "Mode d’emploi",
        "title": "Comment lire ces architectures",
        "bullets": [
            "Chaque schéma décrit un MÉCANISME, pas une règle isolée.",
            "Signal (qui voit) → Jugement (qui décide) → Règle (match event_type) → Preuve (clip + images).",
            "Deux templates qui partagent le même jugement partagent le même schéma.",
            "Fail-closed : pas d’alerte « finale » sans preuves exigées.",
            "Badge catalogue real / partial / beta = vérité produit (un beau schéma ≠ DoD).",
        ],
    },
    {
        "id": "cheat",
        "kind": "table",
        "eyebrow": "Reconnaissance rapide",
        "title": "Quelle question → quel pattern ?",
        "rows": [
            ("Juge habitacle via Gemini yes/no", "A2 Cabin VLM"),
            ("Mesure vitesse km/h", "A3 Measure"),
            ("Feu + véhicule", "A4 Dual-signal"),
            ("Visage / liste surveillance (« liste noire »)", "A1 Face fusion"),
            ("Plaque / OCR / blocklist", "A8 Plate"),
            ("Zone enter / loiter / breach / stopped", "A5 Geometry"),
            ("Ligne / sens interdit", "A6 Line"),
            ("Seuil foule / comptage", "A7 Aggregate"),
            ("Objet abandonné / retiré", "A9 Objects"),
            ("Plusieurs events dans une fenêtre", "A10 Composite"),
        ],
    },
]

ARCHS = [
    {
        "id": "a1",
        "eyebrow": "A1 · Archetype face",
        "title": "Identité faciale — fusion à priorité",
        "blurb": "Crop Frigate → 3 votes (Frigate Face, InsightFace watchlist, Gemini refs). Priorité Frigate > InsightFace > Gemini. « Liste noire » UI = face_watchlist → face_watchlist_match.",
        "same": "tpl-face-watchlist · tpl-watchlist-match · tpl-face-detected · tpl-unknown-face",
        "combine": "tpl-identity-correlation (face + plaque)",
        "evidence": "Clip 6s + scene full + face bbox · fail-closed sur face · identity_votes / winner",
        "debug": "Watchlist + photos enroll + Frigate face ON. Si face_enqueued=0 → bridge, pas Gemini.",
        "mermaid": """%%{init: {'theme':'dark','themeVariables': {'primaryColor':'#1a2f2a','primaryTextColor':'#e8f0ec','primaryBorderColor':'#4a9b84','lineColor':'#8fb8a8','background':'#0e1416','mainBkg':'#1a2f2a','fontFamily':'arial'}}}%%
flowchart LR
  F[Frigate person] --> C[crop JPEG]
  C --> V1[Frigate Face]
  C --> V2[InsightFace]
  C --> V3[Gemini refs]
  V1 --> Fuse[fusion priorité]
  V2 --> Fuse
  V3 --> Fuse
  Fuse --> M[face_watchlist_match]
  Fuse --> U[face_unknown]
  Fuse --> D[face_detected]
  M --> R[rules-engine]
  U --> R
  D --> R
  R --> E[clip + scene + face]
""",
    },
    {
        "id": "a2",
        "eyebrow": "A2 · Archetype cabin",
        "title": "Cabin Gemini — jugement yes / no",
        "blurb": "Zone cabin (seatbelt / phone) → crop vehicle_bbox → file VLM → Gemini. Sans violation claire → cabin_ignored (fail-closed), pas d’alerte.",
        "same": "tpl-seatbelt · tpl-phone-driving",
        "combine": "Ne pas confondre avec A4 feu ou A3 vitesse",
        "evidence": "Clip 6s + scene + subject",
        "debug": "spatial seatbelt chaud → cabin_enqueued>0 → vlm_completed. emitted=0 + rejected=0 ≈ « pas de violation ».",
        "mermaid": """%%{init: {'theme':'dark','themeVariables': {'primaryColor':'#1a2f2a','primaryTextColor':'#e8f0ec','primaryBorderColor':'#4a9b84','lineColor':'#8fb8a8','background':'#0e1416','mainBkg':'#1a2f2a'}}}%%
flowchart LR
  F[Frigate car zone cabin] --> C[crop vehicle_bbox]
  C --> Q[VLM queue]
  Q --> G[Gemini judge]
  G -->|violation| E1[seatbelt / phone event]
  G -->|no unclear| X[cabin_ignored]
  E1 --> R[rules-engine]
  R --> Ev[clip + scene + subject]
""",
    },
    {
        "id": "a3",
        "eyebrow": "A3 · Archetype measure",
        "title": "Vitesse — pont Frigate speed",
        "blurb": "Track + zone de mesure → speed_kmh vs seuil → speeding. Preuve scène + sujet (+ plaque si pipeline plaque).",
        "same": "tpl-speeding · tpl-speeding-premium",
        "combine": "tpl-traffic-pipeline (vitesse + feu + plaque)",
        "evidence": "Clip 6s + scene + subject + plate?",
        "debug": "zone speed_measurement en spatial + bridge speed ON (speed_emitted vs below_limit).",
        "mermaid": """%%{init: {'theme':'dark','themeVariables': {'primaryColor':'#1a2f2a','primaryTextColor':'#e8f0ec','primaryBorderColor':'#4a9b84','lineColor':'#8fb8a8','background':'#0e1416','mainBkg':'#1a2f2a'}}}%%
flowchart LR
  F[Frigate track + zone] --> S[speed_kmh]
  S -->|above| Evt[speeding]
  S -->|below| Drop[below_limit]
  Evt --> R[rules-engine]
  R --> Ev[clip + scene + subject]
""",
    },
    {
        "id": "a4",
        "eyebrow": "A4 · Archetype dual_signal",
        "title": "Feu rouge — géométrie + état feu",
        "blurb": "Véhicule ∩ zone ET feu rouge (HSV / vote). Sinon skipped_not_red. Plus exigeant qu’une simple présence.",
        "same": "tpl-red-light",
        "combine": "traffic-pipeline avec A3 + A8",
        "evidence": "Clip 6s + scene + subject + plate",
        "debug": "Distinguer pas de voiture / feu pas rouge / refus VLM.",
        "mermaid": """%%{init: {'theme':'dark','themeVariables': {'primaryColor':'#1a2f2a','primaryTextColor':'#e8f0ec','primaryBorderColor':'#4a9b84','lineColor':'#8fb8a8','background':'#0e1416','mainBkg':'#1a2f2a'}}}%%
flowchart LR
  F[Frigate car zone] --> H[HSV light state]
  F --> C[crop]
  C --> G[Gemini vote?]
  H --> J[jugement dual]
  G --> J
  J -->|red+cross| Evt[red_light_violation]
  J -->|not red| Skip[skipped_not_red]
  Evt --> R[rules-engine]
""",
    },
    {
        "id": "a5",
        "eyebrow": "A5 · Archetype geometry",
        "title": "Géométrie de zone",
        "blurb": "Présence, loitering/dwell, périmètre, stopped, absence. Jugement spatial-temporel, VLM non obligatoire. Zones via ZoneEditor uniquement.",
        "same": "zone-presence · absence · loitering · perimeter · dwell · person/vehicle_stopped",
        "combine": "composites intrusion + schedule (A10)",
        "evidence": "Clip 6s + scene + subject bbox",
        "debug": "Zone dans Frigate (cv_zone_*) + spatial AI avant d’accuser la règle.",
        "mermaid": """%%{init: {'theme':'dark','themeVariables': {'primaryColor':'#1a2f2a','primaryTextColor':'#e8f0ec','primaryBorderColor':'#4a9b84','lineColor':'#8fb8a8','background':'#0e1416','mainBkg':'#1a2f2a'}}}%%
flowchart LR
  F[Frigate object] --> G[geometry]
  G --> P[presence / enter]
  G --> L[loitering / dwell]
  G --> B[perimeter]
  G --> S[stopped]
  G --> A[absence]
  P --> R[rules-engine]
  L --> R
  B --> R
  S --> R
  A --> R
""",
    },
    {
        "id": "a6",
        "eyebrow": "A6 · Archetype line",
        "title": "Ligne — croisement & sens",
        "blurb": "Ligne (pas zone) : line_cross / sens interdit. Base du comptage directionnel.",
        "same": "tpl-line-cross · bidir · forbidden",
        "combine": "A7 agrégats sur crossings",
        "evidence": "Clip 6s + scene + subject",
        "debug": "Ligne active DB + spatial/Frigate. Ligne seule ≠ gate cabin.",
        "mermaid": """%%{init: {'theme':'dark','themeVariables': {'primaryColor':'#1a2f2a','primaryTextColor':'#e8f0ec','primaryBorderColor':'#4a9b84','lineColor':'#8fb8a8','background':'#0e1416','mainBkg':'#1a2f2a'}}}%%
flowchart LR
  F[Frigate track] --> L[line geometry]
  L -->|cross| Evt[line_cross]
  L -->|forbidden| Forb[line_cross_forbidden]
  Evt --> R[rules-engine]
  Forb --> R
""",
    },
    {
        "id": "a7",
        "eyebrow": "A7 · Archetype aggregate",
        "title": "Agrégats — seuils & densité",
        "blurb": "On agrège (count / densité / fenêtre). Seuil franchi → event. Foule, files, parking.",
        "same": "vehicle-count · crowd-count · crowd-density",
        "combine": "observation N/OR",
        "evidence": "Clip 6s + scene",
        "debug": "Vérifier fenêtre + seuil bindings — souvent un faux bug IA.",
        "mermaid": """%%{init: {'theme':'dark','themeVariables': {'primaryColor':'#1a2f2a','primaryTextColor':'#e8f0ec','primaryBorderColor':'#4a9b84','lineColor':'#8fb8a8','background':'#0e1416','mainBkg':'#1a2f2a'}}}%%
flowchart LR
  F[Frigate objects] --> C[count / density]
  C -->|threshold| Evt[count / crowd event]
  Evt --> R[rules-engine]
""",
    },
    {
        "id": "a8",
        "eyebrow": "A8 · Archetype plate",
        "title": "Plaque — OCR fusion + listes",
        "blurb": "Crop → PaddleOCR (+ Gemini OCR) → fusion → plate_detected puis block / allow / unknown / repeat. Listes plaque ≠ watchlist visage.",
        "same": "plate-detected · blocked · whitelist · unknown · repeat",
        "combine": "A1 face · A3 speed · A4 red light",
        "evidence": "Clip 6s + scene + plate crop",
        "debug": "OCR health + crop lisible + liste active.",
        "mermaid": """%%{init: {'theme':'dark','themeVariables': {'primaryColor':'#1a2f2a','primaryTextColor':'#e8f0ec','primaryBorderColor':'#4a9b84','lineColor':'#8fb8a8','background':'#0e1416','mainBkg':'#1a2f2a'}}}%%
flowchart LR
  F[Frigate vehicle] --> C[plate crop]
  C --> O1[PaddleOCR]
  C --> O2[Gemini OCR]
  O1 --> Fuse[ocr_fusion]
  O2 --> Fuse
  Fuse --> Det[plate_detected]
  Det --> Bl[blocked / unknown / repeat]
  Det --> R[rules-engine]
  Bl --> R
""",
    },
    {
        "id": "a9",
        "eyebrow": "A9 · Archetype objects",
        "title": "Objets — abandon / retrait / disparition",
        "blurb": "Géométrie + temps sur objet (hybrid). Différent du loitering personne (A5).",
        "same": "abandoned · object-removed · object-disappeared",
        "combine": "atomes séparés (pas de theft-composite catalogue)",
        "evidence": "Clip 6s + scene + object bbox",
        "debug": "Durée min abandon/disparition dans bindings.",
        "mermaid": """%%{init: {'theme':'dark','themeVariables': {'primaryColor':'#1a2f2a','primaryTextColor':'#e8f0ec','primaryBorderColor':'#4a9b84','lineColor':'#8fb8a8','background':'#0e1416','mainBkg':'#1a2f2a'}}}%%
flowchart LR
  F[object hybrid] --> G[geometry + time]
  G --> Ab[abandoned]
  G --> Rm[removed]
  G --> Di[disappeared]
  Ab --> R[rules-engine]
  Rm --> R
  Di --> R
""",
    },
    {
        "id": "a10",
        "eyebrow": "A10 · Archetype composite",
        "title": "Composite — séquence d’atomes",
        "blurb": "Pas un nouveau détecteur. Fenêtre sur atoms déjà émis. Preuve = union. Atom missing → suppress fail-closed.",
        "same": "identity-correlation · traffic-pipeline · industrial-intrusion · observation N/OR",
        "combine": "réutilise A1–A9",
        "evidence": "Union des packages atomes",
        "debug": "Valider chaque atom isolément avant le composite.",
        "mermaid": """%%{init: {'theme':'dark','themeVariables': {'primaryColor':'#1a2f2a','primaryTextColor':'#e8f0ec','primaryBorderColor':'#4a9b84','lineColor':'#8fb8a8','background':'#0e1416','mainBkg':'#1a2f2a'}}}%%
flowchart LR
  A1[atom A + evidence] --> W[window]
  A2[atom B + evidence] --> W
  W -->|all OK| Comp[composite]
  W -->|missing| Fail[suppress]
  Comp --> Ev[union preuves]
""",
    },
]

PRODUCT = [
    {
        "id": "b1",
        "eyebrow": "B1 · Runtime",
        "title": "Démarrage Start-CiteVision",
        "blurb": "Launcher → WSL start-full-stack STRICT → infra → API → IA → rules → business readiness → UI → watchdogs. Readiness FAIL = Start non-zéro.",
        "same": "Start-CiteVision.ps1 · start-full-stack.sh",
        "combine": "B2 watchdogs post-start",
        "evidence": "N/A (boot)",
        "debug": "Messages produit streamés — pas de chemins miroir labo côté client.",
        "mermaid": """%%{init: {'theme':'dark','themeVariables': {'primaryColor':'#1a2f2a','primaryTextColor':'#e8f0ec','primaryBorderColor':'#4a9b84','lineColor':'#8fb8a8','background':'#0e1416','mainBkg':'#1a2f2a'}}}%%
flowchart TD
  PS[Start PS1] --> WSL[WSL runtime]
  WSL --> Env[env LIVE_108]
  Env --> Dock[compose infra]
  Dock --> API[backend]
  API --> AI[ai-engine]
  AI --> RE[rules-engine]
  RE --> BR[business readiness]
  BR --> UI[frontend]
  BR -->|FAIL| Abort[exit non-0]
""",
    },
    {
        "id": "b2",
        "eyebrow": "B2 · Resilience",
        "title": "Auto-heal + readiness métier",
        "blurb": "Ports/containers + état métier (spatial, rules, Frigate zones, go2rtc). Probe → heal → re-probe.",
        "same": "watch-infra-ports · watch-business-readiness · service-heal",
        "combine": "B1 Start STRICT",
        "evidence": "N/A",
        "debug": "Redis down ≠ zone_count=0 — heals différents.",
        "mermaid": """%%{init: {'theme':'dark','themeVariables': {'primaryColor':'#1a2f2a','primaryTextColor':'#e8f0ec','primaryBorderColor':'#4a9b84','lineColor':'#8fb8a8','background':'#0e1416','mainBkg':'#1a2f2a'}}}%%
flowchart TD
  P[probe] -->|OK| Idle[OK]
  P -->|KO| H[heal]
  H --> RP[re-probe]
  RP -->|OK| Idle
  RP -->|KO| Fail[FAIL / WARN]
  W[watchdogs] --> P
""",
    },
    {
        "id": "b3",
        "eyebrow": "B3 · Sorties",
        "title": "Alerte → webhook / mail",
        "blurb": "Event → rules match → alert DB → evidence async → routing. Pas de finale si evidence missing exigée.",
        "same": "routing · webhook · mail premium",
        "combine": "B6 preuves",
        "evidence": "complete avant sortie « finale »",
        "debug": "Tester webhook sur alerte complete, pas sur log MQTT seul.",
        "mermaid": """%%{init: {'theme':'dark','themeVariables': {'primaryColor':'#1a2f2a','primaryTextColor':'#e8f0ec','primaryBorderColor':'#4a9b84','lineColor':'#8fb8a8','background':'#0e1416','mainBkg':'#1a2f2a'}}}%%
flowchart LR
  Ev[event] --> RE[rules match]
  RE --> Al[alert]
  Al --> Evd[evidence]
  Evd -->|complete| Out[routing]
  Out --> WH[webhook]
  Out --> Mail[mail]
  Evd -->|missing| Hold[hold]
""",
    },
    {
        "id": "b4",
        "eyebrow": "B4 · Onboarding",
        "title": "Ajout d’une caméra",
        "blurb": "UI → cameras DB → Frigate cv_<uuid> → StartCamera IA + spatial. Zones via ZoneEditor seulement.",
        "same": "onboard UI · orchestrator SyncNow",
        "combine": "B1 ingest LIVE · A* rules",
        "evidence": "N/A jusqu’à première alerte",
        "debug": "Live View ≠ ingest IA. Vérifier spatial après zone/règle.",
        "mermaid": """%%{init: {'theme':'dark','themeVariables': {'primaryColor':'#1a2f2a','primaryTextColor':'#e8f0ec','primaryBorderColor':'#4a9b84','lineColor':'#8fb8a8','background':'#0e1416','mainBkg':'#1a2f2a'}}}%%
flowchart TD
  UI[UI onboard] --> DB[(cameras)]
  DB --> Fr[Frigate rebuild]
  DB --> Ing[StartCamera]
  Ing --> AI[AI spatial]
  UI --> Z[ZoneEditor]
  Z --> RS[resync-spatial]
  RS --> AI
""",
    },
    {
        "id": "b5",
        "eyebrow": "B5 · Infra",
        "title": "Orchestration containers",
        "blurb": "Compose WSL : Postgres, Redis, MQTT, MinIO, go2rtc, Frigate, OCR. API/AI/rules/frontend = process host + watchdogs.",
        "same": "docker compose · dockerd natif WSL",
        "combine": "B2 heal ports",
        "evidence": "N/A",
        "debug": "Ports host morts après hibernation → heal infra d’abord.",
        "mermaid": """%%{init: {'theme':'dark','themeVariables': {'primaryColor':'#1a2f2a','primaryTextColor':'#e8f0ec','primaryBorderColor':'#4a9b84','lineColor':'#8fb8a8','background':'#0e1416','mainBkg':'#1a2f2a'}}}%%
flowchart TD
  C[compose] --> PG[(postgres)]
  C --> R[(redis)]
  C --> M[mqtt]
  C --> S3[(minio)]
  C --> G2[go2rtc]
  C --> Fr[frigate]
  API[backend] --> PG
  AI[ai-engine] --> M
  AI --> Fr
""",
    },
    {
        "id": "b6",
        "eyebrow": "B6 · Données",
        "title": "Preuves — Postgres + MinIO",
        "blurb": "Events/alerts en DB. Binaires (clip, JPEG) dans MinIO. Composer async, jamais de preuve fabriquée. Retention purge démo.",
        "same": "evidence composer · MinIO bucket · retention",
        "combine": "B3 sorties · validate_rule",
        "evidence": "complete | missing (+ cause)",
        "debug": "DoD = alerte UI + fichiers MinIO + mail — pas un log MQTT.",
        "mermaid": """%%{init: {'theme':'dark','themeVariables': {'primaryColor':'#1a2f2a','primaryTextColor':'#e8f0ec','primaryBorderColor':'#4a9b84','lineColor':'#8fb8a8','background':'#0e1416','mainBkg':'#1a2f2a'}}}%%
flowchart LR
  Det[détection] --> Ev[(events)]
  Ev --> Al[(alerts)]
  Al --> Comp[composer]
  Comp --> Clip[clip 6s]
  Comp --> Img[images]
  Clip --> S3[(MinIO)]
  Img --> S3
  Comp --> Meta[evidence_status]
""",
    },
]


def slide_title(prs, item):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_textbox(slide, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.2), item["title"], size=48, bold=True, color=ACCENT, font="Georgia")
    add_textbox(slide, Inches(0.9), Inches(3.3), Inches(11), Inches(0.6), item["subtitle"], size=22, bold=False, color=INK, font="Calibri")
    box = slide.shapes.add_textbox(Inches(0.9), Inches(4.2), Inches(10.5), Inches(2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(item["body"].split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(8)
        run = p.add_run()
        run.text = line
        set_run(run, size=16, color=MUTED)
    add_textbox(slide, Inches(0.9), Inches(6.9), Inches(8), Inches(0.3), "CitéVision · Architectures", size=11, color=TEAL)


def slide_bullets(prs, item):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.35), Inches(11), Inches(0.35), item["eyebrow"], size=12, bold=True, color=TEAL)
    add_textbox(slide, Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.7), item["title"], size=32, bold=True, color=INK, font="Georgia")
    box = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.2), Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(item["bullets"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i:
            p.space_before = Pt(14)
        run = p.add_run()
        run.text = "▸  " + b
        set_run(run, size=18, color=MUTED if i else INK)


def slide_table(prs, item):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.3), item["eyebrow"], size=12, bold=True, color=TEAL)
    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11.5), Inches(0.6), item["title"], size=28, bold=True, color=INK, font="Georgia")
    rows = item["rows"]
    table = slide.shapes.add_table(len(rows) + 1, 2, Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.5)).table
    table.columns[0].width = Inches(7.8)
    table.columns[1].width = Inches(3.9)
    for ci, h in enumerate(["Si la règle…", "Pattern"]):
        cell = table.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                set_run(r, size=12, bold=True, color=ACCENT)
        fill_solid(cell, BG2) if False else None
    for ri, (a, b) in enumerate(rows, start=1):
        table.cell(ri, 0).text = a
        table.cell(ri, 1).text = b
        for ci in (0, 1):
            for p in table.cell(ri, ci).text_frame.paragraphs:
                for r in p.runs:
                    set_run(r, size=12, color=INK if ci == 1 else MUTED, bold=(ci == 1))


def slide_arch(prs, item):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_textbox(slide, Inches(0.55), Inches(0.22), Inches(12), Inches(0.28), item["eyebrow"], size=11, bold=True, color=ACCENT)
    add_textbox(slide, Inches(0.55), Inches(0.48), Inches(12.2), Inches(0.5), item["title"], size=26, bold=True, color=INK, font="Georgia")
    add_textbox(slide, Inches(0.55), Inches(1.0), Inches(12.2), Inches(0.75), item["blurb"], size=13, color=MUTED)

    png = mermaid_to_png(item["mermaid"], item["id"])
    if png:
        # diagram panel
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.85), Inches(8.3), Inches(3.55))
        fill_solid(panel, RGBColor(0x0E, 0x14, 0x16))
        panel.line.color.rgb = RGBColor(0x2A, 0x3A, 0x36)
        slide.shapes.add_picture(str(png), Inches(0.7), Inches(2.05), width=Inches(7.9))
    else:
        add_textbox(slide, Inches(0.7), Inches(2.5), Inches(7.5), Inches(2), "(Diagramme indisponible hors-ligne — voir index.html)", size=14, color=MUTED)

    # side cards
    def side_card(top, label, text, color=TEAL):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0), top, Inches(3.85), Inches(1.05))
        fill_solid(card, BG2)
        card.line.color.rgb = RGBColor(0x2A, 0x3A, 0x36)
        add_textbox(slide, Inches(9.15), top + Inches(0.1), Inches(3.5), Inches(0.25), label, size=10, bold=True, color=color)
        add_textbox(slide, Inches(9.15), top + Inches(0.35), Inches(3.55), Inches(0.65), text, size=11, color=MUTED)

    side_card(Inches(1.85), "MÊME FONCTIONNEMENT", item["same"])
    side_card(Inches(3.05), "COMBINE AVEC", item["combine"], ACCENT)
    side_card(Inches(4.25), "PREUVES", item["evidence"], OK)

    # debug bar
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.6), Inches(12.35), Inches(1.35))
    fill_solid(bar, RGBColor(0x14, 0x22, 0x20))
    bar.line.color.rgb = TEAL
    add_textbox(slide, Inches(0.75), Inches(5.75), Inches(11.8), Inches(0.3), "DEBUG PRAGMATIQUE", size=11, bold=True, color=OK)
    add_textbox(slide, Inches(0.75), Inches(6.1), Inches(11.8), Inches(0.7), item["debug"], size=14, color=INK)


def slide_section(prs, eyebrow, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_textbox(slide, Inches(0.9), Inches(2.6), Inches(11), Inches(0.4), eyebrow, size=14, bold=True, color=TEAL)
    add_textbox(slide, Inches(0.9), Inches(3.1), Inches(11.5), Inches(0.8), title, size=36, bold=True, color=INK, font="Georgia")
    add_textbox(slide, Inches(0.9), Inches(4.0), Inches(10), Inches(1), subtitle, size=16, color=MUTED)


def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_textbox(slide, Inches(0.9), Inches(2.4), Inches(11), Inches(0.5), "CitéVision", size=40, bold=True, color=ACCENT, font="Georgia")
    add_textbox(
        slide,
        Inches(0.9),
        Inches(3.2),
        Inches(11),
        Inches(1.5),
        "De la caméra à la preuve actionnable.\n\nCompanion web : docs/architectures/index.html\nMême contenu, navigation interactive + diagrammes Mermaid.",
        size=16,
        color=MUTED,
    )


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    print("Fetching diagrams…")
    for a in ARCHS + PRODUCT:
        p = mermaid_to_png(a["mermaid"], a["id"])
        print(f"  {a['id']}: {'OK' if p else 'MISSING'}")

    for item in SLIDES:
        if item["kind"] == "title":
            slide_title(prs, item)
        elif item["kind"] == "bullets":
            slide_bullets(prs, item)
        elif item["kind"] == "table":
            slide_table(prs, item)

    slide_section(prs, "PARTIE A", "Dix architectures de règles", "Les patterns qui couvrent le catalogue CitéVision.")
    for a in ARCHS:
        slide_arch(prs, a)

    slide_section(prs, "PARTIE B", "Six architectures produit", "Start, heal, sorties, caméras, containers, preuves.")
    for a in PRODUCT:
        slide_arch(prs, a)

    slide_closing(prs)
    prs.save(OUT)
    print(f"Wrote {OUT} ({OUT.stat().st_size // 1024} KB, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
